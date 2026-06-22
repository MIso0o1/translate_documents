import asyncio
import httpx
from pptx import Presentation
from pptx.shapes.group import GroupShape
from translate_text import translate_text_with_deepl_async

async def translate_run_async(run, target_lang, client):
    original_text = run.text
    if original_text and original_text.strip():
        original_font_name = run.font.name
        original_font_size = run.font.size
        translated_text = await translate_text_with_deepl_async(original_text, target_lang, client)
        if translated_text:
            run.text = translated_text
            run.font.name = original_font_name
            run.font.size = original_font_size

async def process_shape_async(shape, target_lang, client):
    if isinstance(shape, GroupShape):
        tasks = [process_shape_async(s, target_lang, client) for s in shape.shapes]
        if tasks:
            await asyncio.gather(*tasks)
    elif hasattr(shape, "text_frame"):
        tasks = []
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                tasks.append(translate_run_async(run, target_lang, client))
        if tasks:
            await asyncio.gather(*tasks)

async def translate_pptx_file_async(input_pptx_path, target_lang):
    presentation = Presentation(input_pptx_path)
    async with httpx.AsyncClient() as client:
        tasks = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                tasks.append(process_shape_async(shape, target_lang, client))
        if tasks:
            await asyncio.gather(*tasks)
    return presentation

def translate_pptx_file(input_pptx_path, target_lang):
    return asyncio.run(translate_pptx_file_async(input_pptx_path, target_lang))
